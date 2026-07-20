from __future__ import annotations

import shutil
import subprocess
import tempfile
from pathlib import Path
from typing import Callable, Iterable

from presentation_analysis.config import AnalysisConfig
from presentation_analysis.models import SlideContent
from presentation_analysis.ocr_processor import OCRProcessor
from presentation_analysis.text_utils import (
    count_words,
    merge_native_and_ocr,
    no_space_len,
    split_nonempty_lines,
    unique_preserve_order,
)

LogFn = Callable[[str], None]


def find_ppt_file(root: Path, explicit_path: Path | None = None) -> Path | None:
    if explicit_path is not None:
        candidate = explicit_path if explicit_path.is_absolute() else root / explicit_path
        return candidate if candidate.exists() else None

    for name in ("sample.pptx", "sample.ppt"):
        candidate = root / name
        if candidate.exists():
            return candidate
    return None


def extract_presentation(
    ppt_path: Path,
    config: AnalysisConfig,
    logger: LogFn | None = None,
) -> list[SlideContent]:
    suffix = ppt_path.suffix.lower()
    if suffix == ".pptx":
        return extract_pptx(ppt_path, config, logger=logger)
    if suffix == ".ppt":
        return _extract_legacy_ppt(ppt_path, config, logger=logger)
    raise ValueError(f"지원하지 않는 PPT 형식입니다: {ppt_path.suffix}")


def _extract_legacy_ppt(
    ppt_path: Path,
    config: AnalysisConfig,
    logger: LogFn | None = None,
) -> list[SlideContent]:
    executable = shutil.which("libreoffice") or shutil.which("soffice")
    if executable is None:
        raise RuntimeError(
            ".ppt 파일 분석을 위해 LibreOffice가 필요합니다.\n"
            "sample.ppt 파일을 sample.pptx 형식으로 변환한 후 다시 실행해 주세요."
        )

    with tempfile.TemporaryDirectory(prefix="ppt_convert_") as temp_dir:
        temp_path = Path(temp_dir)
        command = [
            executable,
            "--headless",
            "--convert-to",
            "pptx",
            str(ppt_path),
            "--outdir",
            str(temp_path),
        ]
        if logger:
            logger(f".ppt 파일을 임시 .pptx로 변환합니다: {ppt_path.name}")
        completed = subprocess.run(command, capture_output=True, text=True, encoding="utf-8", errors="replace")
        if completed.returncode != 0:
            raise RuntimeError(f"LibreOffice 변환 실패: {completed.stderr.strip() or completed.stdout.strip()}")

        converted = temp_path / f"{ppt_path.stem}.pptx"
        if not converted.exists():
            candidates = list(temp_path.glob("*.pptx"))
            if not candidates:
                raise RuntimeError("LibreOffice 변환 결과 .pptx 파일을 찾을 수 없습니다.")
            converted = candidates[0]
        return extract_pptx(converted, config, logger=logger)


def extract_pptx(
    pptx_path: Path,
    config: AnalysisConfig,
    logger: LogFn | None = None,
) -> list[SlideContent]:
    try:
        from pptx import Presentation  # type: ignore
        from pptx.enum.shapes import MSO_SHAPE_TYPE  # type: ignore
    except Exception as exc:
        raise RuntimeError("PPTX 분석을 위해 python-pptx 패키지가 필요합니다.") from exc

    ocr = OCRProcessor(config)
    for warning in ocr.warnings:
        if logger:
            logger(f"[OCR] {warning}")

    presentation = Presentation(str(pptx_path))
    slides: list[SlideContent] = []

    for index, slide in enumerate(presentation.slides, start=1):
        text_parts: list[str] = []
        ocr_parts: list[str] = []
        warnings: list[str] = []
        text_box_count = 0
        image_count = 0
        table_count = 0
        table_row_count = 0
        table_column_count = 0
        table_cell_count = 0
        chart_count = 0
        flattened_shapes = list(_iter_shapes(slide.shapes))
        shape_count = len(flattened_shapes)

        title = _extract_title(slide)

        for shape in flattened_shapes:
            if getattr(shape, "has_text_frame", False):
                text = _shape_text(shape)
                if text:
                    text_box_count += 1
                    text_parts.append(text)

            if getattr(shape, "has_table", False):
                table_count += 1
                table = shape.table
                rows = len(table.rows)
                cols = len(table.columns)
                table_row_count += rows
                table_column_count += cols
                table_cell_count += rows * cols
                text_parts.extend(_table_texts(table))

            if getattr(shape, "has_chart", False):
                chart_count += 1

            if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.PICTURE:
                image_count += 1
                text, warning = _ocr_shape(shape, ocr)
                if text:
                    ocr_parts.append(text)
                if warning:
                    warnings.append(f"슬라이드 {index}: {warning}")

        speaker_notes = _extract_notes(slide)
        native_text = "\n".join(unique_preserve_order(text_parts))
        raw_ocr_text = "\n".join(unique_preserve_order(ocr_parts))
        ocr_text, combined_text = merge_native_and_ocr(native_text, raw_ocr_text)

        slide_content = SlideContent(
            index=index,
            title=title,
            native_text=native_text,
            ocr_text=ocr_text,
            combined_text=combined_text,
            speaker_notes=speaker_notes,
            native_text_char_count=len(native_text),
            native_text_no_space_count=no_space_len(native_text),
            ocr_text_char_count=len(ocr_text),
            total_text_char_count=len(combined_text),
            total_text_no_space_count=no_space_len(combined_text),
            word_count=count_words(combined_text),
            line_count=len(split_nonempty_lines(combined_text)),
            text_box_count=text_box_count,
            image_count=image_count,
            table_count=table_count,
            table_row_count=table_row_count,
            table_column_count=table_column_count,
            table_cell_count=table_cell_count,
            chart_count=chart_count,
            shape_count=shape_count,
            extraction_warnings=warnings,
        )
        slides.append(slide_content)
        if logger:
            logger(
                f"[PPT] {index}번 슬라이드 추출 완료: "
                f"텍스트 {slide_content.total_text_char_count}자, 이미지 {image_count}개, 표 {table_count}개"
            )

    return slides


def _iter_shapes(shapes: Iterable) -> Iterable:
    for shape in shapes:
        yield shape
        children = getattr(shape, "shapes", None)
        if children is not None:
            yield from _iter_shapes(children)


def _extract_title(slide) -> str:
    try:
        if slide.shapes.title and slide.shapes.title.text.strip():
            return slide.shapes.title.text.strip()
    except Exception:
        pass

    for shape in slide.shapes:
        text = _shape_text(shape)
        if text:
            return split_nonempty_lines(text)[0][:80]
    return "제목 없음"


def _shape_text(shape) -> str:
    try:
        text = shape.text
    except Exception:
        return ""
    return "\n".join(unique_preserve_order(split_nonempty_lines(text)))


def _table_texts(table) -> list[str]:
    texts: list[str] = []
    for row in table.rows:
        for cell in row.cells:
            text = "\n".join(unique_preserve_order(split_nonempty_lines(cell.text)))
            if text:
                texts.append(text)
    return texts


def _ocr_shape(shape, ocr: OCRProcessor) -> tuple[str, str | None]:
    try:
        blob = shape.image.blob
    except Exception as exc:
        return "", f"이미지 데이터를 읽을 수 없습니다: {exc}"
    return ocr.extract_text_from_image_blob(blob)


def _extract_notes(slide) -> str:
    try:
        if not slide.has_notes_slide:
            return ""
        notes_frame = slide.notes_slide.notes_text_frame
        return "\n".join(unique_preserve_order(split_nonempty_lines(notes_frame.text)))
    except Exception:
        return ""
