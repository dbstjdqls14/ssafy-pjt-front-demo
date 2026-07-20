from __future__ import annotations

import tempfile
import unittest
from pathlib import Path
from xml.etree import ElementTree as ET

from PIL import Image
from pptx import Presentation

from presentation_analysis.config import AnalysisConfig
from presentation_analysis.content_density import apply_density
from presentation_analysis.models import SlideContent, STTSegment
from presentation_analysis.ppt_extractor import extract_presentation
from presentation_analysis.presentation_evaluator import evaluate_slides
from presentation_analysis.report_writer import write_reports
from presentation_analysis.runner import run_presentation_analysis
from presentation_analysis.slide_aligner import align_slides
from presentation_analysis.stt_processor import analyze_stt_segments
from presentation_analysis.time_recommender import build_recommendations


class PresentationAnalysisTests(unittest.TestCase):
    def test_pptx_text_table_and_ocr_fallback(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            temp_path = Path(temp_dir)
            pptx_path = temp_path / "table_image.pptx"
            image_path = temp_path / "image.png"
            Image.new("RGB", (160, 100), "white").save(image_path)

            presentation = Presentation()
            slide = presentation.slides.add_slide(presentation.slide_layouts[5])
            slide.shapes.title.text = "표 테스트"
            slide.shapes.add_textbox(100000, 900000, 4000000, 600000).text = "본문 텍스트"
            table_shape = slide.shapes.add_table(2, 2, 100000, 1600000, 4000000, 1000000)
            table_shape.table.cell(0, 0).text = "항목"
            table_shape.table.cell(0, 1).text = "값"
            table_shape.table.cell(1, 0).text = "정확도"
            table_shape.table.cell(1, 1).text = "95%"
            slide.shapes.add_picture(str(image_path), 100000, 2800000, width=1400000, height=900000)
            presentation.save(pptx_path)

            slides = extract_presentation(pptx_path, AnalysisConfig())
            self.assertEqual(len(slides), 1)
            self.assertIn("본문 텍스트", slides[0].native_text)
            self.assertIn("정확도", slides[0].native_text)
            self.assertEqual(slides[0].table_count, 1)
            self.assertEqual(slides[0].table_cell_count, 4)
            self.assertEqual(slides[0].image_count, 1)
            self.assertTrue(any("OCR 비활성화" in warning for warning in slides[0].extraction_warnings))

    def test_missing_ppt_message_does_not_crash(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            messages: list[str] = []
            result = run_presentation_analysis(
                root=Path(temp_dir),
                analyze_ppt_only=True,
                logger=messages.append,
            )
            self.assertIsNone(result)
            self.assertIn("루트 폴더에서 sample.pptx 또는 sample.ppt 파일을 찾을 수 없습니다.", messages)

    def test_stt_events_alignment_long_slide_and_xml_output(self) -> None:
        config = AnalysisConfig()
        slides = [
            _slide(1, "표지", "서비스 소개"),
            _slide(2, "문제 정의", "문제 설명"),
        ]
        apply_density(slides, config)
        stt = analyze_stt_segments(
            [
                STTSegment(0.0, 8.0, "서비스를 소개합니다"),
                STTSegment(10.0, 86.0, "문제를 반복해서 자세히 설명합니다 문제를 반복해서 자세히 설명합니다"),
            ],
            config,
            audio_duration=90.0,
        )

        with tempfile.TemporaryDirectory() as temp_dir:
            temp_path = Path(temp_dir)
            events_path = temp_path / "slide_events.csv"
            events_path.write_text("slide_index,start_time,end_time\n1,0,10\n2,10,86\n", encoding="utf-8")
            alignments, _note = align_slides(slides, stt, config, events_path)
            recommendations, recommended_total, buffer_seconds = build_recommendations(slides, stt.chars_per_second, config)
            evaluations = evaluate_slides(slides, alignments, recommendations, stt, config)

            self.assertEqual(alignments[1].actual_duration, 76.0)
            self.assertTrue(evaluations[1].is_too_long)
            self.assertTrue(evaluations[1].reasons)

            reports = write_reports(
                temp_path / "sample.pptx",
                slides,
                stt,
                alignments,
                recommendations,
                evaluations,
                None,
                recommended_total,
                buffer_seconds,
                "slide_events.csv 기반 실제 전환 시간",
                temp_path / "logs",
            )
            self.assertEqual(ET.parse(reports.ppt_content_xml).getroot().tag, "presentation")
            self.assertEqual(ET.parse(reports.analysis_xml).getroot().tag, "presentationAnalysis")
            self.assertIn("문제 정의", Path(reports.analysis_txt).read_text(encoding="utf-8"))


def _slide(index: int, title: str, text: str) -> SlideContent:
    return SlideContent(
        index=index,
        title=title,
        native_text=text,
        ocr_text="",
        combined_text=text,
        speaker_notes="",
        native_text_char_count=len(text),
        native_text_no_space_count=len(text.replace(" ", "")),
        ocr_text_char_count=0,
        total_text_char_count=len(text),
        total_text_no_space_count=len(text.replace(" ", "")),
        word_count=len(text.split()),
        line_count=1,
        text_box_count=1,
        image_count=0,
        table_count=0,
        table_row_count=0,
        table_column_count=0,
        table_cell_count=0,
        chart_count=0,
        shape_count=1,
    )


if __name__ == "__main__":
    unittest.main()
